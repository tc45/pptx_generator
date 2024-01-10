import io

from pptx import Presentation
from pptx.util import Pt, Inches
import markdown
import re


def markdown_to_pptx(markdown_text):
    pres = Presentation()
    parser = markdown.Markdown()

    # Convert Markdown to HTML
    html = parser.convert(markdown_text)

    # Regex Patterns
    h1_pattern = r'<h1>(.*?)</h1>'
    h2_pattern = r'<h2>(.*?)</h2>'
    ul_pattern = r'<ul>(.*?)</ul>'
    li_pattern = r'<li>(.*?)</li>'
    paragraph_pattern = r'<p>(.*?)<\/p>'
    note_pattern = r'<!--(.*?)-->'

    # Process H1 - Title Slide
    h1_matches = re.findall(h1_pattern, html, re.DOTALL)
    h1_positions = re.search(h1_pattern, html, re.DOTALL)
    h2_positions = re.search(h2_pattern, html, re.DOTALL)
    paragraph_positions = re.search(paragraph_pattern, html, re.DOTALL)
    paragraph_matches = re.findall(paragraph_pattern, html, re.DOTALL)
    notes_positions = re.search(note_pattern, html, re.DOTALL)
    notes_matches = re.findall(paragraph_pattern, html, re.DOTALL)

    if h1_matches:
        slide_layout = pres.slide_layouts[0]  # Title Slide
        slide = pres.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = h1_matches[0]
        if paragraph_matches is not None:
            h1_span = h1_positions.span()
            p_span = paragraph_positions.span()
            h2_span = h2_positions.span()
            # Check if the paragraph match is between H1 and H2.  If so, add as subtitle.
            if p_span[0] > h1_span[1] and p_span[1] > h2_span[0]:
                paragraph_match = re.findall(paragraph_pattern, html, re.DOTALL)
                subtitle.text = paragraph_match[0]

    # Process H2 and Bullet Points
    for h2_match in re.finditer(h2_pattern, html, re.DOTALL):
        html_sliced = html[h2_match.start():]

        slide_layout = pres.slide_layouts[1]  # Title and Content
        slide = pres.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = h2_match.group(1)

        # Extracting bullet points and sub bullet points
        start, end = h2_match.span()
        ul_content = re.search(ul_pattern, html[end:], re.DOTALL)
        if ul_content:
            content = slide.placeholders[1]
            for li_match in re.finditer(li_pattern, ul_content.group(1), re.DOTALL):
                p = content.text_frame.add_paragraph()
                p.text = re.sub(r'<.*?>', '', li_match.group(1))  # Remove HTML tags
                p.level = 0
                p.font.size = Pt(18)
                print(p.text)

        # Adding notes to the slide
        notes = slide.notes_slide.notes_text_frame
        note_content = re.search(note_pattern, html_sliced, re.DOTALL)

        if note_content:
            notes.text = note_content.group(1).strip()
            print(note_content.group(1))
    return pres


def check_file_extension(file_name):
    extension = file_name[-4:].lower()
    if file_name[-4:].lower() != "pptx" and file_name[-3:].lower() != "ppt":
        if file_name[-1:] != ".":
            extension = "." + "pptx"
        else:
            extension = "pptx"
        file_name = file_name + extension
        return file_name


def save_presentation(markdown, file_path, format=""):
    filename = check_file_extension(file_path)
    print("File is " + filename)
    pres = markdown_to_pptx(markdown)
    if format == "":
        pres.save(filename)
        return pres
    elif format == "bytes":
        pptx_io = io.BytesIO()
        pres.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io, filename

# Example usage
markdown_text = """# TEST2 Title
This is the main title page.
Subtitle 2

## Slide 1
- Bullet 1-1
- Bullet 2-1
    - Sub Bullet 1
<!-- Note for Slide 1 21321`31241234-->

## Slide 2
- Bullet 2-1
- Bullet 2-2
<!-- Note for Slide 2 126532651412-->"""


if __name__ == '__main__':
    save_presentation(markdown_text, 'output.pptx2')
