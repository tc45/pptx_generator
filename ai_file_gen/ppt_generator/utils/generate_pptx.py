import io
from pptx import Presentation
from pptx.util import Pt
import markdown
import re


# Example usage
markdown_text = """
# Introduction to Network Automation
test paragraph
- Understanding the basics of network automation
- The importance of network automation in modern networks
- Key benefits: efficiency, accuracy, and scalability
## Ansible
- Overview of Ansible
- Key features and advantages
- Ansible Playbooks and Roles
## Python
- Why Python for network automation?
- Common libraries: Netmiko, Paramiko, NAPALM
## Identifying Automation Opportunities
- Analyzing repetitive tasks
- Prioritizing tasks for automation
- Planning and executing automation projects
## Basic Pythonic Syntax for Automation
- Python data types and structures
- Control structures: loops and conditionals
- Writing reusable functions
## Case Study: Using Netmiko
- Introduction to Netmiko
- Connecting to routers and switches
- Automating common network tasks
## Best Practices in Network Automation
- Writing maintainable and scalable code
- Error handling and logging
- Security considerations
"""


def markdown_to_pptx(markdown_text):
    """
    Convert Markdown text to PowerPoint.
    Args:
        markdown_text (str): Input text in markdown.
    Returns:
        Presentation: An instance of Presentation class (from python-pptx) containing all slides and configuration.
    """
    # Create a new presentation, Markdown object
    pres = Presentation()
    parser = markdown.Markdown()

    # Convert Markdown to HTML
    html = parser.convert(markdown_text)

    # Initialize regex patterns
    h1_pattern = r'<h1>(.*?)</h1>'
    h2_pattern = r'<h2>(.*?)</h2>'
    ul_pattern = r'<ul>(.*?)</ul>'
    li_pattern = r'<li>(.*?)</li>'
    paragraph_pattern = r'<p>(.*?)<\/p>'
    note_pattern = r'<!--(.*?)-->'

    # Search for markdown tags in HTML
    h1_matches = re.findall(h1_pattern, html, re.DOTALL)
    h1_positions = re.search(h1_pattern, html, re.DOTALL)
    h2_positions = re.search(h2_pattern, html, re.DOTALL)
    paragraph_positions = re.search(paragraph_pattern, html, re.DOTALL)
    paragraph_matches = re.findall(paragraph_pattern, html, re.DOTALL)
    notes_positions = re.search(note_pattern, html, re.DOTALL)
    notes_matches = re.findall(note_pattern, html, re.DOTALL)

    # Convert H1 markdown tag to PPT title slide
    if h1_matches:
        slide_layout = pres.slide_layouts[0]  # Title Slide
        slide = pres.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        # Set title text
        title.text = h1_matches[0]

        # Check for paragraphs between H1 and H2 to add as subtitle.
        if paragraph_matches:
            h1_span = h1_positions.span()
            p_span = paragraph_positions.span()
            h2_span = h2_positions.span()

            if p_span[0] > h1_span[1] and p_span[1] < h2_span[0]:
                paragraph_match = re.findall(paragraph_pattern, html, re.DOTALL)
                subtitle.text = paragraph_match[0]

    # Convert each H2 markdown tag to new PPT slide
    for h2_match in re.finditer(h2_pattern, html, re.DOTALL):
        html_sliced = html[h2_match.start():]
        slide_layout = pres.slide_layouts[1]  # Title and Content
        slide = pres.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = h2_match.group(1)

        # Get start, end spans
        start, end = h2_match.span()
        ul_content = re.search(ul_pattern, html[end:], re.DOTALL)

        # Deal with unordered lists as bullet points
        if ul_content:
            content = slide.placeholders[1]
            for li_match in re.finditer(li_pattern, ul_content.group(1), re.DOTALL):
                p = content.text_frame.add_paragraph()
                p.text = re.sub(r'<.*?>', '', li_match.group(1))  # Remove HTML tags
                p.level = 0
                p.font.size = Pt(18)
                print(p.text)

        # Adding notes to the slide if available
        notes = slide.notes_slide.notes_text_frame
        note_content = re.search(note_pattern, html_sliced, re.DOTALL)
        if note_content:
            notes.text = note_content.group(1).strip()
            print(note_content.group(1))

    return pres


def check_file_extension(file_name):
    """
    Check the file extension of the user-specified filename. If it ends in either .ppt or .pptx nothing is done. A .pptx extension is added otherwise.

    Args:
    filename (str): Input filename.

    Returns:
    str: Output filename. If the original one didn't have an extension or had a different one, the new filename has a .pptx extension.
    """
    extension = file_name[-4:].lower()
    # check and correct extension if not ending with .ppt or .pptx
    if extension != ".pptx" and extension != ".ppt":
        if file_name[-1:] != ".":
            extension = ".pptx"
        else:
            extension = "pptx"
        file_name += extension

    return file_name


def save_presentation(markdown, file_path, format=""):
    """
    Checks for the file_path. Sets the default path to 'presentation.pptx' if no path is provided. Generates PowerPoint presentation from input markdown text. Presentation is saved as .pptx file or as bytes depending upon the format input.

    Args:
    markdown (str): markdown text to convert to PowerPoint presentation.
    file_path (str): Path where the generated ppt will be stored.
    format (str, optional): Specifies if bytes format for file is required.Defaults to "".

    Returns:
    Presentation/BytesIO, str: Presentation instance or BytesIO object of the generated ppt, filename.
    """
    # Default path
    if file_path is None: file_path = "presentation.pptx"

    # Check and correct file extension for pptx
    filename = check_file_extension(file_path)

    # Generate PowerPoint from Markdown
    pres = markdown_to_pptx(markdown)

    # Save result to byte array or .pptx file
    if format == "":
        pres.save(filename)
        return pres
    elif format == "bytes":
        pptx_io = io.BytesIO()
        pres.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io, filename


def test_gpt():
    """
    Connect with the GPT API and send queries to it.

    Returns:
    str: response coming from the api.
    """
    # Import libraries
    from openai import OpenAI, OpenAIError
    from decouple import config

    # Get GPT_API_KEY value from environment variables
    api_key_env = config('GPT_API_KEY')

    # Raise ValueError if value of GPT_API_KEY is not found
    if not api_key_env:
        raise ValueError("GPT API key not found in environment variables")

    # Assign GPT_API_KEY value to api_key
    api_key = api_key_env
    client = OpenAI(api_key=api_key)
    response = None
    try:
        stream = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": "Translate the text to French: 'Hello, how are you?'"}
            ],
            stream=True)

        # Collects replies from OpenAI
        collected_messages = []
        for chunk in stream:
            if chunk.choices[0].delta.content is not None:
                collected_messages.append(chunk.choices[0].delta.content)
        content = "".join(collected_messages)
        return content
    except OpenAIError as e:
        return f"Error: {str(e)}"


if __name__ == '__main__':
    print(test_gpt())
    # Save markdown text as a .pptx file
    save_presentation(markdown_text, 'output.pptx2')