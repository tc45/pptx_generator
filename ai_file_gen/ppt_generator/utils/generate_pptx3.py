import io
import re
import markdown
from pptx import Presentation
from pptx.util import Inches, Pt
import datetime


def save_markdown_text(markdown_text, file_path):
    timestamp = datetime.datetime.now().strftime('%Y%m%d%H%M%S')
    filename = f"{file_path}_markup_{timestamp}.txt"
    with open(filename, 'w') as f:
        f.write(markdown_text)


def extract_html_tags(html, tag):
    """
    Extract contents between given HTML tags.
    :param html: HTML Text.
    :param tag: HTML Tag to extract.
    :return: Content between HTML tags.
    """
    pattern = f'<{tag}>(.*?)</{tag}>'
    return re.findall(pattern, html, re.DOTALL), re.search(pattern, html, re.DOTALL)


def extract_bullets(slide, html):
    """
    Organizes bullet points to a slide from the HTML text.
    :param slide: Current slide to add bullet points.
    :param html: HTML to extract bullet points from.
    """
    # Set UL and LI tag patterns to search for bullet points
    ul_pattern = r'<ul>(.*?)</ul>'
    li_pattern = r'<li>(.*?)</li>'
    ul_content = re.search(ul_pattern, html, re.DOTALL)

    # If unsorted list content is found
    if ul_content:
        # Create a new text box
        left = Inches(2)
        top = Inches(1)
        width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)

        # Access the text_frame property of the textbox
        tf = txBox.text_frame

        # Get list items within unsorted list
        for li_match in re.finditer(li_pattern, ul_content.group(1), re.DOTALL):
            p = tf.add_paragraph()
            # Strip HTML tags from the list item text and add as content
            p.text = re.sub(r'<.*?>', '', li_match.group(1))
            p.level = 0
            p.font.size = Pt(18)


def create_slide(pres, layout, match, html):
    """
    Create a new slide with specifics on layout, extracted content, and handle bullet points.
    :param pres: Presentation.
    :param layout: Layout of the slide.
    :param match: Extracted sections from re.finditer.
    :param html: HTML to extract bullet points from.
    :return: Slide with title and bullet points.
    """
    # Set Slide Layout
    slide_layout = pres.slide_layouts[layout]
    slide = pres.slides.add_slide(slide_layout)

    # Set Section Title
    slide.shapes.title.text = match.group(1)

    # Extract Bullet Points
    extract_bullets(slide, html[match.start():])

    return slide


def strip_empty_lines(text):
    """
    Strips all empty lines from the given text.
    :param str text: the text from which empty lines should be removed.
    :returns: a version of the text that does not contain any empty lines.
    :rtype: str
    """
    lines = text.split("\n")
    non_empty_lines = [line for line in lines if line.strip() != ""]

    return "\n".join(non_empty_lines)

def markdown_to_pptx(markdown_text):
    """
    Convert Markdown text to PowerPoint Presentation.
    :param markdown_text: Markdown text to be converted.
    :return: Presentation file.
    """
    # Strip empty lines from markdown
    markdown_text = strip_empty_lines(markdown_text)
    print("==============MARKDOWN_FINAL============\n\n\n" + markdown_text + "\n\n\n================END=============\n\n")

    # Presentation Initialization
    pres = Presentation()
    parser = markdown.Markdown()
    html = parser.convert(markdown_text)
    print("=============HTML==========\n\n\n" + html + "\n\n\===========END==========\n\n")

    # Parsing H1, H2, H3 Sections to Create Slides
    for match in re.finditer("<h1>(.*?)</h1>", html, re.DOTALL):
        print("Found Title slide" + match.group(1))
        create_slide(pres, 0, match, html)
    for match in re.finditer("<h2>(.*?)</h2>", html, re.DOTALL):
        create_slide(pres, 1, match, html)
    for match in re.finditer("<h3>(.*?)</h3>", html, re.DOTALL):
        create_slide(pres, 2, match, html)  # Assuming layout for H3 is 2

    return pres


def check_file_extension(file_name):
    """
    Check PPTX file extension availability and append if not present.
    :param file_name:
    :return: File name with PPTX extension.
    """
    if not file_name.endswith(('.ppt', '.pptx')):
        file_name += '.pptx'
    return file_name


def save_presentation(markdown_text, file_path, format=""):
    """
    Save the presentation.
    :param markdown_text: Markdown text to be converted.
    :param file_path: File path to store the converted PPTX.
    :param format: Format to return the PPTX (default '', optional 'bytes').
    :return: If format is 'bytes' return pptx_io and filename else None.
    """
    # Save the markdown text locally
    save_markdown_text(markdown_text, file_path)
    print('Markdown Details: \n\n\n\n' + markdown_text)
    if not file_path:
        file_path = 'presentation.pptx'
    filename = check_file_extension(file_path)
    pres = markdown_to_pptx(markdown_text)

    if format == "bytes":
        # Save presentation to BytesIO stream and return
        pptx_io = io.BytesIO()
        pres.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io, filename
    else:
        # Save Presentation as PPTX file
        pres.save(filename)


# Example usage
markdown_text = """
# Introduction to Network Automation
- Understanding the basics of network automation
- The importance of network automation in modern networks
- Key benefits: efficiency, accuracy, and scalability

## Benefits of using automation
- Improve productivity
- Reduce human performance error
- Ensure consistency
- Scalablility

## Differences Between a VAR and an End User of Automation
- Focus
  - VAR: Often focuses on integrating and customizing solutions for a variety of clients.
  - End User: Focuses on using solutions to meet specific internal business needs.
- Scope of Work
  - VAR: Works with a range of products and services to provide comprehensive solutions.
  - End User: Typically concentrated on specific tools and platforms relevant to their business.
- Objective
  - VAR: Aims to add value to products/services and resell to clients with markup.
  - End User: Looks to optimize and streamline internal processes and operations.
- Skills and Expertise
  - VAR: Requires broad knowledge across multiple products and how they integrate.
  - End User: More likely to have in-depth knowledge of particular tools relevant to their domain.
- Revenue Generation
  - VAR: Generates revenue by selling enhanced products/services.
  - End User: Leverages automation to reduce costs and increase efficiency, indirectly impacting revenue.
- Client Interaction
  - VAR: Regularly engages with different clients, requiring strong customer relations skills.
  - End User: Primarily internal interactions, focusing on internal team dynamics and requirements.

## Identifying Automation Opportunities
- Analyzing repetitive tasks in network management
- Prioritizing tasks for automation
- Planning and executing automation projects

## Tools and Technologies for Network Automation
- Overview of Python
  - Key features and advantages
  - Libraries, vendor support, industry standard
- Overview of Ansible
  - Key features and advantages
  - Ansible Playbooks and Roles

## Using AI to Enhance Automation
- Overview
  - Integration of AI technologies is revolutionizing automation by adding intelligence and adaptability.
- ChatGPT
  - Utilized for natural language processing tasks.Enhances automation by enabling conversational interfaces and understanding user intents.
- GitHub Copilot
  - AI-powered code completion tool.
- AI Assistants**
  - Virtual assistants powered by AI to automate routine tasks (schedule meetings, send reminders, write powerpoints)
- Predictive Analytics
  - AI-driven analytics to predict trends and automate decision-making processes.
- Machine Learning in Network Management
  - Automating network configuration and maintenance based on learned network patterns and behaviors.
- AI in Cybersecurity
  - Automated threat detection and response.
  
## Benefits of using AI for automation
  - Increased efficiency and accuracy.
  - Reduced manual effort and time spent on routine tasks.
  - Enhanced capability to handle complex scenarios and data analysis.


## Python Basics for Automation
- Why Python for network automation?
- Python data types and structures
- Control structures: loops and conditionals
- Writing reusable functions

## Development Environments
- Overview of Integrated Development Environments (IDEs)
  - Popular IDEs for Python and network automation (e.g., PyCharm, Visual Studio Code)
- Benefits of using cloud-based IDEs
  - Examples of cloud IDEs (e.g., AWS Cloud9, Gitpod)

## Common Python Network Tools
- Netmiko
  - Automating configuration changes, retrieving information, and executing commands on network devices.
- Paramiko
  - Automating tasks that require secure connections to network devices or servers.
- NAPALM (Network Automation and Programmability Abstraction Layer with Multivendor support)
  - Gathering data and automating configurations across multiple network vendors.
- Ansible
  - Managing configurations, deploying policies, and orchestrating complex network operations.
- Scapy
  - Creating, sending, capturing, and analyzing network packets. Useful in network research and security testing.
- PySNMP
  - Performing SNMP operations like querying and setting MIBs on network devices.
- ExaBGP
  - Manipulating BGP routes, useful for tasks like network testing, load balancing, and DDoS mitigation.
- Libnmap
  - Network discovery and security auditing.
- Python-iptables
  - Managing firewall rules and packet filtering.
- Requests
  - Interacting with REST APIs, often used in network and security automation for communicating with network controllers, security appliances, and cloud services.


## Common Python Tools
- Pandas for data manipulation and analysis
- Requests for making HTTP requests and interacting with APIs


## Ansible Fundamentals
- What is Ansible? 
  - Overview and origin
  - Ansible vs. other automation tools
- Core Concepts of Ansible
  - Playbooks, Roles, Tasks, and Modules
  - Inventory files and Configuration management
- Installation and Setup
  - Installing Ansible on various operating systems
  - Setting up a basic environment for testing

## Working with Ansible
- Writing Your First Playbook
  - Basic structure and syntax
  - Running a simple playbook
- Understanding Ansible Modules
  - Commonly used modules in network automation
  - Creating custom modules
- Advanced Ansible Features
  - Using Ansible roles for reusable content
  - Templating with Jinja2
  - Ansible Vault for managing secrets

## Ansible for Network Automation
- Automating Network Devices with Ansible
  - Supported devices and connection methods
  - Example tasks for network configuration and management
- Best Practices in Ansible for Networking
  - Structuring playbooks for network tasks
  - Handling errors and ensuring idempotency

## Version Control with Git
- Introduction to Git and version control
- Setting up and managing Git repositories
- Cloud-based Git Platforms: GitHub and GitLab
  - Features and benefits of cloud-based version control
- Git Workflow
  - Common Git workflows (e.g., feature branch, Gitflow)
  - Best practices for collaborative development

## Continuous Integration and Continuous Deployment (CI/CD)
- Principles of CI/CD in network automation
- Tools and platforms for CI/CD (e.g., Jenkins, Travis CI)
- Automating network configuration and deployment processes

## Netmiko Case Study: Using Netmiko: Connecting to routers and switches
- In this exmaple we will explore connecting to routers and switches using SSH

## Netmiko Case Study: Using Netmiko: Collect data and store
- In this example we will look at collecting and storing data into a local file

## Netmiko Case Study: Using Netmiko: Parse collected data into JSON
- In this example we will look at how to use the textfsm capabilities to parse collected data

## Ansible Case Study Review: Automating Network Configuration Changes
- Background: The need for automating network changes in a mid-sized company
- Challenge: Reducing manual effort and errors in network configuration
- Solution: Implementing Ansible for automated deployment and updates
- Process:
  - Identifying repetitive network tasks
  - Writing Ansible Playbooks for automating these tasks
  - Testing and rolling out in a staged approach
- Results:
  - Time savings and reduced manual errors
  - Improved consistency in network configurations
  - Enhanced ability to quickly respond to network changes
  
## Ansible Case Study Example: Automating network configuration changes
  - In this section we will look at an example of how to use ansible to automate network changes

## Ansible Case Study: Streamlining Data Center Operations
- Background: A large enterprise with a complex data center infrastructure
- Challenge: Managing a vast number of servers and network devices, ensuring consistency and compliance with industry standards
- Solution: Implementing Ansible for configuration management and orchestration
- Process:
  - Mapping out all the network devices and server configurations
  - Creating Ansible Playbooks for routine tasks like patches, updates, and configuration changes
  - Integrating Ansible with existing CI/CD pipelines for seamless operations
- Results:
  - Significant reduction in configuration errors and manual intervention
  - Faster deployment of updates and new configurations across the data center
  - Enhanced security and compliance with automated checks and configurations

## Best Practices in Network Automation
- Writing maintainable and scalable code
- Error handling and logging in automation scripts
- Security considerations in network automation

## Expand Your Knowledge
- Online Courses and Certifications
  - Recommendations for advanced courses in network automation (e.g., specific Udemy or Coursera courses)
  - Industry certifications that cover network automation (e.g., Cisco DevNet, Juniper Networks Automation)
- Books and Reading Materials
  - Essential books for deeper understanding (e.g., "Network Programmability and Automation" by Jason Edelman, Scott S. Lowe, and Matt Oswalt)
  - Online resources and blogs for staying updated

## Join the Community
- Networking forums and online communities (e.g., Stack Overflow, Reddit's r/networking, and r/devops)
- Local meetups and user groups focusing on network automation
- Conferences and workshops (e.g., Cisco Live, DevNet Create)

## Practice and Hands-on Experience
- Setting up a personal lab for practice
  - Using virtual labs and simulators (e.g., GNS3, Cisco's VIRL)
- Contributing to open-source projects related to network automation
- Experimenting with APIs of commonly used networking tools

## Stay Updated with Industry Trends
- Following key influencers and thought leaders in network automation on social media
- Subscribing to newsletters and podcasts in the field
- Regularly reviewing and experimenting with new tools and technologies

## Call to Action
- Start small: Automate a simple task in your network environment.
- Share your learnings with peers and colleagues.
- Challenge yourself with a personal project that utilizes network automation.
- Keep learning: Dedicate regular time to update your skills and knowledge.

"""

if __name__ == '__main__':
    save_presentation(markdown_text, 'output.pptx')
