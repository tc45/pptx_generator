import re
from pptx import Presentation
from pptx.util import Pt
import markdown

# Regex Patterns
h_patterns = [
    (r'<h1>(.*?)</h1>', 0),
    (r'<h2>(.*?)</h2>', 1),
    (r'<h3>(.*?)</h3>', 1)
]
ul_pattern = r'<ul>(.*?)</ul>'
li_pattern = r'<li>(.*?)</li>'
note_pattern = r'<!--(.*?)-->'
html = ""


def handle_title(pres, title):
    slide_layout = pres.slide_layouts[0]  # Title Slide
    slide = pres.slides.add_slide(slide_layout)
    slide_title = slide.shapes.title
    slide_title.text = title[0]


def handle_headings(pres, headings, level):
    for heading in headings:
        slide_layout = pres.slide_layouts[1]  # Content Slide
        slide = pres.slides.add_slide(slide_layout)
        slide_title = slide.shapes.title
        slide_title.text = heading
        check_for_bullets(pres, slide, 0)


def handle_notes(pres, notes):
    for note in notes:
        slide_layout = pres.slide_layouts[2]  # Notes Slide
        slide = pres.slides.add_slide(slide_layout)
        notes_placeholder = slide.notes_slide.notes_text_frame
        notes_placeholder.text = note.strip()


def check_for_bullets(pres, slide, level):
    ul_content = re.search(ul_pattern, html, re.DOTALL)
    if ul_content:
        content = slide.placeholders[1]
        for li_match in re.finditer(li_pattern, ul_content.group(1), re.DOTALL):
            p = content.text_frame.add_paragraph()
            p.text = re.sub(r'<.*?>', '', li_match.group(1))  # Remove HTML tags
            p.level = level
            p.font.size = Pt(18)


def markdown_to_pptx(markdown_text):
    pres = Presentation()
    parser = markdown.Markdown()

    # Convert Markdown to HTML
    html = parser.convert(markdown_text)

    # Handling the H1, H2, H3 tags
    for pattern, level in h_patterns:
        matches = re.findall(pattern, html, re.DOTALL)
        if matches:
            handle_headings(pres, matches, level)

    # Handling the notes
    notes = re.findall(note_pattern, html, re.DOTALL)
    if notes:
        handle_notes(pres, notes)

    return pres